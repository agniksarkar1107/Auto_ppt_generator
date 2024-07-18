
from dotenv import load_dotenv
load_dotenv()

import os
import streamlit as st
import google.generativeai as genai
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re


from pptx import Presentation



os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)



def get_gemini_response(prompt):
    model = genai.GenerativeModel('gemini-1.5-flash')
    response = model.generate_content(prompt)
    return response

sub_titles = []
def refine_subtopics(sub_topics, sub_titles):
    for sub_topic in sub_topics:
        sub_titles.append(sub_topic[3:].replace('"',""))
    return sub_titles

content = []
def content_generation(sub_titles):
    for i in sub_titles:
        prompt = f"Generate a content of {i} for the presentation topic of {topic} in a structured manner in short paragraphs highlighting the important words in around 100 words"
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        content.append(response.text)
    return content

def clean_text(text):
    # Remove extra whitespaces and newlines
    cleaned_text = re.sub('\s+', ' ', text).strip()

    # Remove markdown-style bullet points, asterisks, and numeric bullet points
    cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)

    # Remove extra spaces before and after colons
    cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)

    # Remove extra spaces before and after hyphens
    cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)

    return cleaned_text


def split_sentences(text):
    # Split the text into sentences using regular expression
    sentences = re.split(r'(?<=\.)\s+', text)

    # Capitalize the first letter of each sentence
    sentences = [sentence.capitalize() for sentence in sentences]

    return sentences
def replace_and_capitalize(text):
    # Define a function to replace and capitalize the text between colons
    def replace_and_capitalize_colon(match):
        return match.group(1) + match.group(2).capitalize() + match.group(3)

    # Use regular expression to find and replace text between colons
    result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)', replace_and_capitalize_colon, text)

    return result


final_content = []
def refine_final_content(content):
    for i in content:
        cleaned_text = clean_text(i)
        sentences = split_sentences(cleaned_text)
        final_content.append(sentences)
    print("final content ready....")
    return final_content

powerpoint = Presentation()

def slide_maker(powerpoint, topic,sub_titles, final_content):
    title_slide_layout = powerpoint.slide_layouts[0]
    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]
    content.text = "Created by Ethan's ppt generator"
    for i in range(len(sub_titles)):
        bulletLayout = powerpoint.slide_layouts[1]
        secondSlide = powerpoint.slides.add_slide(bulletLayout)
        # accessing the attributes of shapes
        myShapes = secondSlide.shapes
        titleShape = myShapes.title
        bodyShape = myShapes.placeholders[1]
        titleShape.text = sub_titles[i]
        titleShape.text_frame.paragraphs[0].font.size = Pt(24)
        titleShape.text_frame.paragraphs[0].font.bold = True
        tFrame = bodyShape.text_frame
        print("Topic Generated")
        for point in final_content[i]:
            point = re.sub(r':[^:]+:', ':', point)
            point = replace_and_capitalize(point)
            p = tFrame.add_paragraph()
            p.text = point
            p.font.size = Pt(15)
            p.space_after = Pt(10)
    return powerpoint
def download_button(file_path,topic):
    # Read the content of the PPT file
    with open(file_path, "rb") as file:
        ppt_content = file.read()

    # Create a download button for the PPT file
    st.download_button(
        label="Download PowerPoint",
        data=ppt_content,
        file_name=f"{topic}.pptx",  # Change the file name as needed
        key="ppt_download_button"
    )

#model = genai.GenerativeModel('gemini-pro')
st.set_page_config(page_title="Ethan's Presentation Maker")

st.header("Ethan's Presentation Maker")
topic=st.text_input("Input Prompt: ",key="input")
no_of_slide=st.text_input("Enter Number Of Slide: ",key="slide")

submit=st.button("Generate")
if submit:
    prompt =f"Generate {no_of_slide} very short simple sub-headings for a presentation only  on the topic of {topic}"
    response = get_gemini_response(prompt)
    print("Topic Generated")
    sub_topics = response.text.split("\n")
    sub_titles = refine_subtopics(sub_topics, sub_titles)
    print("Sub Titles")
    content = content_generation(sub_titles)
    print("content Generated")
    final_content = refine_final_content(content)
    #cleaned_text = clean_text(content[0])
    #sentences = split_sentences(cleaned_text)
    print("final content ready")
    powerpoint = slide_maker(powerpoint,topic, sub_titles, final_content)
    
    
    binary_output = BytesIO()
    powerpoint.save(binary_output) 

    st.download_button(label = 'Download ppw',
                   data = binary_output.getvalue(),
                   file_name = 'my_power.pptx')

    
